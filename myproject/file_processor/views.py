import os
import re
import nltk
from django.shortcuts import render
from .forms import ChatForm
from .models import Document
from docx import Document as DocxDocument
from pptx import Presentation
from PyPDF2 import PdfReader  # Для обработки PDF
from transformers import pipeline, DistilBertTokenizer, DistilBertForQuestionAnswering
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize
from nltk.stem import WordNetLemmatizer
from langchain_openai import OpenAI
from langchain_core.prompts import PromptTemplate
from transformers import BloomForCausalLM, BloomTokenizerFast

# Загрузка нужных ресурсов NLTK
nltk.download('punkt')
nltk.download('stopwords')
nltk.download('wordnet')

# Инициализация инструментов для очистки и лемматизации
stop_words = set(stopwords.words("russian"))  # Стоп-слова на русском языке
lemmatizer = WordNetLemmatizer()

# Загрузка модели DistilBERT через Hugging Face pipeline
tokenizer_distilbert = DistilBertTokenizer.from_pretrained('distilbert-base-uncased')
model_distilbert = DistilBertForQuestionAnswering.from_pretrained('distilbert-base-uncased')
nlp_distilbert = pipeline("question-answering", model=model_distilbert, tokenizer=tokenizer_distilbert)

# Инициализация модели OpenAI
llm = OpenAI(
    model="gpt-3.5-turbo-instruct",
    temperature=0.7,
    max_retries=2,
    api_key="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"  # Замените на свой API ключ
)

# Инициализация модели Bloom (для генерации текста)
tokenizer_bloom = BloomTokenizerFast.from_pretrained("bigscience/bloom-560m")
bloom_model = BloomForCausalLM.from_pretrained("bigscience/bloom-560m")

# Шаблон для OpenAI
prompt = PromptTemplate.from_template("How to say {input} in {output_language}:\n")


def clean_text(text):
    """Очистка текста от ненужных символов."""
    text = re.sub(r'\s+', ' ', text)  # Заменяем несколько пробелов на один
    text = re.sub(r'[^\w\s]', '', text)  # Удаляем все небуквенные символы
    return text.strip()


def preprocess_text(text):
    """Предобработка текста: очистка, токенизация, лемматизация, стоп-слова."""
    cleaned_text = clean_text(text)
    cleaned_text = cleaned_text.lower()
    tokens = word_tokenize(cleaned_text)
    tokens = [lemmatizer.lemmatize(word) for word in tokens if word not in stop_words]
    processed_text = " ".join(tokens)
    return processed_text


def extract_text_from_docx(file_path):
    """Извлекает текст из файла .docx."""
    doc = DocxDocument(file_path)
    return '\n'.join([para.text for para in doc.paragraphs])


def extract_text_from_pptx(file_path):
    """Извлекает текст из файла .pptx."""
    prs = Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return '\n'.join(text)


def extract_text_from_pdf(file_path):
    """Извлекает текст из файла .pdf."""
    reader = PdfReader(file_path)
    text = []
    for page in reader.pages:
        text.append(page.extract_text())
    return '\n'.join(text)


def extract_text_from_txt(file_path):
    """Извлекает текст из файла .txt."""
    with open(file_path, 'r', encoding='utf-8') as file:
        return file.read()


# def process_text_with_ai(extracted_text, question, model_name):
#     """Процесс обработки текста и получения ответа от выбранной модели."""
#     if model_name == 'openai':
#         try:
#             chain = prompt | llm
#             response = chain.invoke({
#                 "output_language": "English",  # Язык может быть параметром
#                 "input": question
#             })
#             return response.strip()
#         except Exception as e:
#             return {'error': str(e)}
#
#     elif model_name == 'distilbert':
#         try:
#             result = nlp_distilbert({'context': extracted_text, 'question': question})
#             return result['answer']
#         except Exception as e:
#             return {'error': str(e)}
#
#     elif model_name == 'bloom':
#         try:
#             inputs = tokenizer_bloom(question, return_tensors="pt")
#             output = bloom_model.generate(inputs['input_ids'], max_length=200)
#             return tokenizer_bloom.decode(output[0], skip_special_tokens=True)
#         except Exception as e:
#             return {'error': str(e)}
#
#     else:
#         return "Unsupported model"


def upload_file(request):
    """Обрабатывает загрузку файла и его анализ."""
    ai_responses = []
    if request.method == 'POST':
        form = ChatForm(request.POST, request.FILES)
        if form.is_valid():
            document = form.save()
            file_path = document.file.path
            model_name = form.cleaned_data['model_choice']  # Получаем выбранную модель

            # Извлечение текста в зависимости от типа файла
            if file_path.endswith('.docx'):
                extracted_text = extract_text_from_docx(file_path)
            elif file_path.endswith('.pptx'):
                extracted_text = extract_text_from_pptx(file_path)
            elif file_path.endswith('.pdf'):
                extracted_text = extract_text_from_pdf(file_path)
            elif file_path.endswith('.txt'):
                extracted_text = extract_text_from_txt(file_path)
            else:
                extracted_text = "Unsupported file type."

            # Проверка текста
            if not extracted_text.strip():
                return render(request, 'file_processor/result.html', {'ai_response': "Файл не содержит текста."})

            # Предобработка текста
            processed_text = preprocess_text(extracted_text)

            # Разделение текста на вопросы (например, если они находятся в отдельных строках)
            questions = [line for line in processed_text.split('\n') if line.strip()]

            # Отправка текста в чат
            ai_responses = [{'question': 'Текст из файла', 'answer': extracted_text}]

            # Получение ответов на вопросы
            for question in questions:
                ai_response = process_text_with_ai(processed_text, question, model_name)
                ai_responses.append({'question': question, 'answer': ai_response})

            return render(request, 'file_processor/chat.html', {'form': form, 'ai_responses': ai_responses})

    else:
        form = ChatForm()
    return render(request, 'file_processor/upload.html', {'form': form})


# def chat_view(request):
#     """Чат с пользователем."""
#     ai_response = None
#     model_name = None
#
#     if request.method == 'POST':
#         form = ChatForm(request.POST)
#         if form.is_valid():
#             user_message = form.cleaned_data['user_message']
#             model_name = form.cleaned_data['model_choice']  # Получаем выбранную модель
#
#             # Если пользователь ввел сообщение, получаем ответ от AI
#             if user_message.strip():
#                 ai_response = process_text_with_ai("", user_message, model_name)
#             else:
#                 ai_response = "Пожалуйста, введите текстовое сообщение."
#
#             return render(request, 'file_processor/chat.html', {'form': form, 'ai_response': ai_response})
#
#     else:
#         form = ChatForm()
#
#     return render(request, 'file_processor/chat.html', {'form': form, 'ai_response': ai_response})

def process_text_with_ai(user_message, model_name):
    """Обрабатывает текст и получает ответ от модели."""
    if model_name == 'openai':
        try:
            chain = prompt | llm
            response = chain.invoke({
                "output_language": "English",
                "input": user_message
            })
            return response.strip()
        except Exception as e:
            return {'error': str(e)}

    elif model_name == 'distilbert':
        try:
            result = nlp_distilbert({'context': user_message, 'question': user_message})
            return result['answer']
        except Exception as e:
            return {'error': str(e)}

    elif model_name == 'bloom':
        try:
            inputs = tokenizer_bloom(user_message, return_tensors="pt")
            output = bloom_model.generate(inputs['input_ids'], max_length=200)
            return tokenizer_bloom.decode(output[0], skip_special_tokens=True)
        except Exception as e:
            return {'error': str(e)}

    else:
        return "Unsupported model"

def chat_view(request):
    """Обрабатывает чат с пользователем."""
    ai_response = None
    if request.method == 'POST':
        form = ChatForm(request.POST)
        if form.is_valid():
            user_message = form.cleaned_data['user_message']
            model_name = form.cleaned_data['model_choice']  # Получаем выбранную модель

            # Если пользователь ввел сообщение, получаем ответ от AI
            if user_message.strip():
                ai_response = process_text_with_ai(user_message, model_name)
            else:
                ai_response = "Пожалуйста, введите текстовое сообщение."

            return render(request, 'file_processor/chat.html', {'form': form, 'ai_response': ai_response})

    else:
        form = ChatForm()

    return render(request, 'file_processor/chat.html', {'form': form, 'ai_response': ai_response})

